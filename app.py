# app.py
import streamlit as st
import pandas as pd
import io
import os
from datetime import datetime
from reportlab.lib.pagesizes import A4
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image as RLImage
from reportlab.lib.styles import getSampleStyleSheet
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
import re

# ----------------- CONFIG / SECRETS -----------------
ADMIN_USERNAME = st.secrets.get("admin", {}).get("username", "admin")
ADMIN_PASSWORD = st.secrets.get("admin", {}).get("password", "ete123")

EXCEL_FILE = "ETE_Robotics-Bom-Data-for-softwares-development.xlsx"
LOGO_FILE = "ETE-Robotics-Logo.png"

# ----------------- UTIL: Load BOM Robustly -----------------
@st.cache_data
def load_excel_data(path=EXCEL_FILE):
    if not os.path.exists(path):
        return pd.DataFrame()
    # try header row around row 12 (index 11) or row 13 (index 12)
    for header_idx in (11, 12, 10, 9):
        try:
            df = pd.read_excel(path, header=header_idx, usecols="B:H", engine="openpyxl")
            # normalize column names
            df.columns = [str(c).strip().replace("\n", " ").replace("\r", "") for c in df.columns]
            # map similar column names to canonical
            rename_map = {}
            for col in df.columns:
                low = col.lower()
                if "head" in low:
                    rename_map[col] = "Head"
                elif "description" in low:
                    rename_map[col] = "Description"
                elif "model" in low or "key spec" in low or "model / key" in low:
                    rename_map[col] = "Model/Key Spec"
                elif "unit" in low and "cost" in low:
                    rename_map[col] = "Unit Cost"
                elif "s.no" in low or "s.no." in low or "s.no" in low:
                    rename_map[col] = "S.no"
                elif "qty" == low or "quantity" in low:
                    rename_map[col] = "Qty"
            df = df.rename(columns=rename_map)
            # check minimal columns
            if "Head" in df.columns and "Model/Key Spec" in df.columns:
                df = df.dropna(subset=["Head"])
                # ensure Unit Cost numeric
                if "Unit Cost" in df.columns:
                    df["Unit Cost"] = pd.to_numeric(df["Unit Cost"].astype(str).str.replace(r'[^\d\.\-]', '', regex=True), errors="coerce").fillna(0)
                return df
        except Exception:
            continue
    # fallback: try reading without header and take row 12 as header
    try:
        df = pd.read_excel(path, header=None, engine="openpyxl")
        if df.shape[0] >= 12:
            header = df.iloc[11].astype(str).tolist()
            df2 = pd.read_excel(path, header=12, engine="openpyxl")
            df2.columns = header
            # repeat normalization
            df2.columns = [str(c).strip().replace("\n", " ").replace("\r", "") for c in df2.columns]
            return df2
    except Exception:
        pass
    return pd.DataFrame()

bom_df = load_excel_data()

# ----------------- UTIL: model/spec splitting & dropdown list -----------------
def split_spec_values(cell):
    if not cell or pd.isna(cell):
        return []
    cell = str(cell)
    # split on pipe | or ' I ' or '/' or ';'
    parts = re.split(r'\s*\|\s*|\s+I\s+|/|;|,', cell)
    parts = [p.strip() for p in parts if p.strip()]
    return parts

@st.cache_data
def build_model_options(df):
    options = []
    if df is None or df.empty:
        return options
    col = "Model/Key Spec"
    if col not in df.columns:
        return options
    for val in df[col].dropna().astype(str):
        options.extend(split_spec_values(val))
    return sorted(list(set(options)))

MODEL_OPTIONS = build_model_options(bom_df)

def find_unit_cost_for_model(df, chosen_model):
    if df is None or df.empty or not chosen_model:
        return 0.0
    # try to find any row where 'Model/Key Spec' contains chosen_model
    mask = df["Model/Key Spec"].astype(str).str.contains(re.escape(chosen_model), case=False, na=False)
    if mask.any():
        vals = df.loc[mask, "Unit Cost"].dropna().astype(float)
        if not vals.empty:
            return float(vals.iloc[0])
    return 0.0

# ----------------- UI: Style for breadcrumb -----------------
BREADCRUMB_CSS = """
<style>
.breadcrumbs {
  display:flex;
  justify-content:center;
  gap:12px;
  margin-top:18px;
  margin-bottom:12px;
  font-weight:600;
  font-size:16px;
}
.bc-item {
  padding:8px 14px;
  border-radius:8px;
  background-color:#f1f3f4;
  cursor:pointer;
}
.bc-item.active {
  background:linear-gradient(90deg,#ff8800,#ffb06b);
  color:white;
  box-shadow: 0 2px 8px rgba(0,0,0,0.12);
}
</style>
"""

# ----------------- Helper: PDF generation -----------------
def create_pdf(customer_info, requirements, rfq_checklist, items):
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4, rightMargin=36, leftMargin=36, topMargin=36, bottomMargin=36)
    styles = getSampleStyleSheet()
    story = []

    # logo
    if os.path.exists(LOGO_FILE):
        try:
            story.append(RLImage(LOGO_FILE, width=160, height=60))
        except Exception:
            pass
    story.append(Spacer(1, 6))
    story.append(Paragraph("ETE Robotics Systems Integrator — RFQ Summary", styles['Title']))
    story.append(Spacer(1, 12))

    # Customer Info (two-column table)
    cust_data = [[Paragraph("<b>Customer</b>", styles['Normal']), ""]]
    for k, v in customer_info.items():
        cust_data.append([Paragraph(f"<b>{k}</b>", styles['Normal']), Paragraph(str(v), styles['Normal'])])
    t = Table(cust_data, colWidths=[120, 360])
    t.setStyle(TableStyle([('VALIGN', (0,0), (-1,-1), 'TOP'), ('INNERGRID', (0,0), (-1,-1), 0.25, colors.grey)]))
    story.append(t)
    story.append(Spacer(1, 12))

    # Requirements summary
    story.append(Paragraph("<b>Requirements</b>", styles['Heading3']))
    req_data = []
    for k, v in requirements.items():
        req_data.append([Paragraph(str(k), styles['Normal']), Paragraph(str(v), styles['Normal'])])
    rt = Table(req_data, colWidths=[160, 320])
    rt.setStyle(TableStyle([('VALIGN',(0,0),(-1,-1),'TOP')]))
    story.append(rt)
    story.append(Spacer(1, 12))

    # RFQ checklist
    story.append(Paragraph("<b>RFQ Checklist</b>", styles['Heading3']))
    chk_data = []
    for k, v in rfq_checklist.items():
        chk_data.append([Paragraph(str(k), styles['Normal']), Paragraph(str(v), styles['Normal'])])
    ct = Table(chk_data, colWidths=[160, 320])
    ct.setStyle(TableStyle([('VALIGN',(0,0),(-1,-1),'TOP')]))
    story.append(ct)
    story.append(Spacer(1, 12))

    # Items table (BOM)
    if items:
        story.append(Paragraph("<b>Selected Items</b>", styles['Heading3']))
        table_data = [["S.no", "Head", "Model/Spec", "Qty", "Unit Cost", "Line Cost (INR)"]]
        total = 0.0
        for i, it in enumerate(items, 1):
            line = [it.get("S.no",""), it.get("Head",""), it.get("ModelSpec",""), str(it.get("Qty","")), f"{it.get('UnitCost',0):,.2f}", f"{it.get('LineCost',0):,.2f}"]
            table_data.append(line)
            total += float(it.get("LineCost", 0))
        table_data.append(["", "", "", "", "Total", f"{total:,.2f}"])
        tbl = Table(table_data, colWidths=[40, 140, 160, 50, 80, 90])
        tbl.setStyle(TableStyle([
            ('GRID',(0,0),(-1,-1),0.3,colors.grey),
            ('BACKGROUND',(0,0),(-1,0),colors.HexColor("#f5f5f5")),
            ('ALIGN',(-2,1),(-1,-1),'RIGHT'),
        ]))
        story.append(tbl)
        story.append(Spacer(1, 12))

    doc.build(story)
    buf.seek(0)
    return buf

# ----------------- Helper: PPT generation -----------------
def create_ppt(customer_info, requirements, rfq_checklist, items):
    prs = Presentation()
    # title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "RFQ Summary"
    subtitle = slide.placeholders[1]
    subtitle.text = f"ETE Robotics — {customer_info.get('Customer Name','')}"
    # attempt logo
    if os.path.exists(LOGO_FILE):
        try:
            slide.shapes.add_picture(LOGO_FILE, Inches(6), Inches(0.3), width=Inches(2))
        except Exception:
            pass

    # customer slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Customer & Requirements"
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(4)).text_frame
    for k,v in customer_info.items():
        p = tx.add_paragraph()
        p.text = f"{k}: {v}"
        p.level = 0
    tx.add_paragraph()
    p = tx.add_paragraph()
    p.text = "Requirements:"
    for k,v in requirements.items():
        pp = tx.add_paragraph()
        pp.text = f"{k}: {v}"
        pp.level = 1

    # items slide
    if items:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Selected Items & Costs"
        rows = len(items) + 2
        cols = 5
        left = Inches(0.4)
        top = Inches(1.4)
        width = Inches(9)
        height = Inches(0.8)
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        # headers
        table.cell(0,0).text = "S.no"
        table.cell(0,1).text = "Model/Spec"
        table.cell(0,2).text = "Head"
        table.cell(0,3).text = "Qty"
        table.cell(0,4).text = "Line Cost (INR)"
        r = 1
        total = 0.0
        for it in items:
            table.cell(r,0).text = str(it.get("S.no",""))
            table.cell(r,1).text = str(it.get("ModelSpec",""))
            table.cell(r,2).text = str(it.get("Head",""))
            table.cell(r,3).text = str(it.get("Qty",""))
            table.cell(r,4).text = f"{it.get('LineCost',0):,.2f}"
            total += float(it.get('LineCost',0))
            r += 1
        table.cell(r-1,0).text = ""
        table.cell(r-1,1).text = ""
        table.cell(r-1,2).text = "Total"
        table.cell(r-1,3).text = ""
        table.cell(r-1,4).text = f"{total:,.2f}"

    # final slide - RFQ checklist
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "RFQ Checklist"
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(4)).text_frame
    for k,v in rfq_checklist.items():
        p = tx.add_paragraph()
        p.text = f"{k}: {v}"
    # return bytes
    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out

# ----------------- Application State -----------------
if "step" not in st.session_state:
    st.session_state.step = 1
if "customer_info" not in st.session_state:
    st.session_state.customer_info = {}
if "requirements" not in st.session_state:
    st.session_state.requirements = {}
if "rfq_checklist" not in st.session_state:
    st.session_state.rfq_checklist = {}
if "selected_items" not in st.session_state:
    st.session_state.selected_items = []

# ----------------- Top UI: Title & Breadcrumb -----------------
st.set_page_config(page_title="ETE RFQ Builder", layout="wide")
st.image(LOGO_FILE if os.path.exists(LOGO_FILE) else None, width=220)
st.markdown(BREADCRUMB_CSS, unsafe_allow_html=True)

# Breadcrumb (centered)
col1, col2, col3, col4, col5 = st.columns([1,4,1,4,1])
with col2:
    classes = []
    def bc_click(step_no):
        st.session_state.step = step_no
    for i, label in enumerate(["Step 1: Customer Info", "Step 2: RFQ Checklist", "Step 3: Submit & Generate"], start=1):
        if st.session_state.step == i:
            clicked = st.button(label, key=f"bc_{i}", help="Click to go to this step")
            if clicked:
                st.session_state.step = i
        else:
            clicked = st.button(label, key=f"bc_{i}_off")
            if clicked:
                st.session_state.step = i

st.write("---")

# ----------------- STEP 1 -----------------
if st.session_state.step == 1:
    st.header("Step 1 — Customer Info & Requirements")
    # Use two-column layout for fields
    with st.form("step1_form"):
        c1, c2 = st.columns(2)
        with c1:
            rfq_ref = st.text_input("RFQ Reference", value=st.session_state.customer_info.get("RFQ Reference", f"RFQ/ETE/{datetime.now().year}-{datetime.now().strftime('%m%d%H%M')}"))
            cust_name = st.text_input("Customer Name", value=st.session_state.customer_info.get("Customer Name",""))
            contact_no = st.text_input("Contact No.", value=st.session_state.customer_info.get("Contact No.",""))
            email = st.text_input("Email ID", value=st.session_state.customer_info.get("Email ID",""))
            location = st.text_input("Location", value=st.session_state.customer_info.get("Location",""))
        with c2:
            date_val = st.date_input("Date", value=st.session_state.customer_info.get("Date", datetime.today()))
            application = st.selectbox("Application", ["Robotic","SPM","Testing","Conveyor","Plant facility Transfer","Modification","Service","Std item Supply"], index=0)
            eq_type = st.selectbox("Type of Equipment", ["Hydraulic","Pneumatic","Servo","Other","Conveyor (Belt)","Conveyor (Slat)","Conveyor (Roller)"], index=0)
            product_info = st.text_input("Product Information", value=st.session_state.customer_info.get("Product Information",""))
            new_mod = st.selectbox("New / Modification", ["New","Modification"], index=0)

        col_btn1, col_btn2, col_btn3 = st.columns([1,1,1])
        with col_btn1:
            save = st.form_submit_button("Save")
        with col_btn2:
            next_btn = st.form_submit_button("Save & Next")
        with col_btn3:
            reset = st.form_submit_button("Reset")
    # handle buttons
    if save or next_btn:
        st.session_state.customer_info = {
            "RFQ Reference": rfq_ref,
            "Customer Name": cust_name,
            "Contact No.": contact_no,
            "Email ID": email,
            "Date": str(date_val),
            "Location": location,
            "Application": application,
            "Type of Equipment": eq_type,
            "Product Information": product_info,
            "New/Modification": new_mod
        }
        st.success("Step 1 data saved.")
        if next_btn:
            st.session_state.step = 2
            st.experimental_rerun()
    if reset:
        st.session_state.customer_info = {}
        st.experimental_rerun()

# ----------------- STEP 2 -----------------
elif st.session_state.step == 2:
    st.header("Step 2 — RFQ Checklist & Bill of Quantity")
    with st.form("step2_form"):
        # A) General Information (two column)
        st.subheader("A) General Information")
        g1, g2 = st.columns(2)
        with g1:
            project_desc = st.text_input("Project Description", value=st.session_state.rfq_checklist.get("Project Description",""))
            proposal_no = st.text_input("Proposal No.", value=st.session_state.rfq_checklist.get("Proposal No","P-001"))
            assigned_to = st.text_input("Assigned To", value=st.session_state.rfq_checklist.get("Assigned To",""))
        with g2:
            customer_repeat = st.text_input("Customer (repeat)", value=st.session_state.rfq_checklist.get("Customer", st.session_state.customer_info.get("Customer Name","")))
            chk_date = st.date_input("Date", value=datetime.today())

        # B) Project Summary
        st.subheader("B) Project Summary")
        s1, s2 = st.columns(2)
        with s1:
            ps_type = st.selectbox("Type of Equipment (summary)", ["Hydraulic","Pneumatic","Servo","Other"], index=0)
            ps_application = st.selectbox("Application (summary)", ["Robotic","SPM","Testing","Conveyor"], index=0)
            cycle_time = st.text_input("Cycle Time", value=st.session_state.rfq_checklist.get("Cycle Time",""))
        with s2:
            no_of_variants = st.number_input("No. of Variants", min_value=1, value=int(st.session_state.rfq_checklist.get("No. of Variants",1)))
            delivery_time = st.text_input("Delivery Time", value=st.session_state.rfq_checklist.get("Delivery Time",""))
            i_and_c_location = st.text_input("I&C Location", value=st.session_state.rfq_checklist.get("I&C Location",""))

        st.subheader("C) Concept Layout")
        concept_layout = st.text_area("Concept Layout (notes)", value=st.session_state.rfq_checklist.get("Concept Layout",""))

        st.subheader("D) Process Details / Key Features")
        # Put a compact multi-line input
        key_features = st.text_area("Key Features / Process Details", value=st.session_state.rfq_checklist.get("Key Features",""), help="Describe key features separated by commas or new lines.")

        st.subheader("E) Bill of Quantity")
        # BOM item selector: model from MODEL_OPTIONS, qty, Add
        bom_col1, bom_col2, bom_col3 = st.columns([2,1,1])
        with bom_col1:
            chosen_model = st.selectbox("Select Model / Key Spec (from BOM)", ["-- select --"] + MODEL_OPTIONS, index=0)
        with bom_col2:
            chosen_qty = st.number_input("Qty", min_value=1, value=1)
        with bom_col3:
            add_item_btn = st.form_submit_button("Add Item to BOM")

        # show selected_items
        st.markdown("**Current selected items**")
        if st.session_state.selected_items:
            df_items = pd.DataFrame(st.session_state.selected_items)
            st.dataframe(df_items[["ModelSpec","Head","Qty","UnitCost","LineCost"]], use_container_width=True)
        else:
            st.info("No BOM items added yet. Use the selector above to add items from BOM.")

        # buttons at bottom
        col_a, col_b, col_c = st.columns([1,1,1])
        with col_a:
            save2 = st.form_submit_button("Save")
        with col_b:
            next2 = st.form_submit_button("Save & Next")
        with col_c:
            back2 = st.form_submit_button("Back to Step 1")

    # handle adding
    if add_item_btn and chosen_model and chosen_model != "-- select --":
        unit_cost = find_unit_cost_for_model(bom_df, chosen_model)
        line_cost = unit_cost * chosen_qty
        # find a head value from bom_df if possible
        head_val = ""
        mask = bom_df["Model/Key Spec"].astype(str).str.contains(re.escape(chosen_model), case=False, na=False) if not bom_df.empty else []
        if getattr(mask, "any", lambda : False)():
            try:
                head_val = bom_df.loc[mask, "Head"].iloc[0]
            except Exception:
                head_val = ""
        item = {"S.no": len(st.session_state.selected_items)+1,
                "ModelSpec": chosen_model,
                "Head": head_val,
                "Qty": int(chosen_qty),
                "UnitCost": float(unit_cost),
                "LineCost": float(line_cost)}
        st.session_state.selected_items.append(item)
        st.success(f"Added {chosen_model} x {chosen_qty}")

    # handle save/next/back
    if save2:
        st.session_state.rfq_checklist = {
            "Project Description": project_desc,
            "Proposal No": proposal_no,
            "Assigned To": assigned_to,
            "Customer": customer_repeat,
            "Date": str(chk_date),
            "Project Summary": {
                "Type": ps_type,
                "Application": ps_application,
                "Cycle Time": cycle_time,
                "Variants": no_of_variants,
                "Delivery Time": delivery_time,
                "I&C Location": i_and_c_location
            },
            "Concept Layout": concept_layout,
            "Key Features": key_features
        }
        st.success("Step 2 saved.")
    if next2:
        # save then go to step 3
        st.session_state.rfq_checklist = {
            "Project Description": project_desc,
            "Proposal No": proposal_no,
            "Assigned To": assigned_to,
            "Customer": customer_repeat,
            "Date": str(chk_date),
            "Project Summary": {
                "Type": ps_type,
                "Application": ps_application,
                "Cycle Time": cycle_time,
                "Variants": no_of_variants,
                "Delivery Time": delivery_time,
                "I&C Location": i_and_c_location
            },
            "Concept Layout": concept_layout,
            "Key Features": key_features
        }
        st.session_state.step = 3
        st.experimental_rerun()
    if back2:
        st.session_state.step = 1
        st.experimental_rerun()

# ----------------- STEP 3 -----------------
elif st.session_state.step == 3:
    st.header("Step 3 — Review, Generate PDF & PPT (No Email)")
    st.subheader("Review entered data")
    st.markdown("**Customer Info**")
    st.json(st.session_state.customer_info)
    st.markdown("**Requirements**")
    st.json(st.session_state.requirements)
    st.markdown("**RFQ Checklist**")
    st.json(st.session_state.rfq_checklist)
    st.markdown("**Selected BOM Items**")
    if st.session_state.selected_items:
        st.dataframe(pd.DataFrame(st.session_state.selected_items)[["ModelSpec","Head","Qty","UnitCost","LineCost"]], use_container_width=True)
    else:
        st.info("No items selected in BOM. Go back to Step 2 and add items.")

    # compute totals
    total = sum([it.get("LineCost",0) for it in st.session_state.selected_items]) if st.session_state.selected_items else 0.0
    st.markdown(f"### Total Estimated Cost: **₹ {total:,.2f}**")

    col_left, col_right = st.columns([1,1])
    with col_left:
        if st.button("Back to Step 2"):
            st.session_state.step = 2
            st.experimental_rerun()
    with col_right:
        if st.button("Generate PDF & PPT"):
            pdf_bytes = create_pdf(st.session_state.customer_info, st.session_state.requirements, st.session_state.rfq_checklist, st.session_state.selected_items)
            ppt_bytes = create_ppt(st.session_state.customer_info, st.session_state.requirements, st.session_state.rfq_checklist, st.session_state.selected_items)
            # expose downloads
            st.success("Generated PDF and PPT. Download links below.")
            st.download_button("Download PDF", data=pdf_bytes, file_name=f"{st.session_state.customer_info.get('RFQ Reference','RFQ')}.pdf", mime="application/pdf")
            st.download_button("Download PPTX", data=ppt_bytes, file_name=f"{st.session_state.customer_info.get('RFQ Reference','RFQ')}.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
            # store in session for later (optional)
            st.session_state.last_pdf = pdf_bytes.getvalue()
            st.session_state.last_pptx = ppt_bytes.getvalue()
            # no email sending as per request
